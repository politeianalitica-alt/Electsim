"""
Office Parser — Bloque 9.

Parsea DOCX, XLSX, PPTX. Degrada sin romper si faltan dependencias.
"""
from __future__ import annotations

import logging
from datetime import datetime, timezone
from pathlib import Path
from typing import TYPE_CHECKING

if TYPE_CHECKING:
    from etl.sources.documents.schemas import ParsedDocument, ExtractedTable

logger = logging.getLogger(__name__)

_DOCX_OK = False
_OPENPYXL_OK = False
_PANDAS_OK = False

try:
    import docx  # python-docx
    _DOCX_OK = True
except ImportError:
    pass

try:
    import openpyxl  # noqa: F401
    _OPENPYXL_OK = True
except ImportError:
    pass

try:
    import pandas  # noqa: F401
    _PANDAS_OK = True
except ImportError:
    pass


def parse_docx(path: str | Path) -> "ParsedDocument":
    """Parsea un archivo DOCX con python-docx."""
    from etl.sources.documents.schemas import ParsedDocument
    path = Path(path)
    document_id = path.stem
    now = datetime.now(timezone.utc)

    if not _DOCX_OK:
        return ParsedDocument(
            document_id=document_id,
            parser_used="docx_fallback",
            warnings=["python-docx no instalado. Instala: pip install python-docx"],
            parsed_at=now,
        )

    try:
        import docx as _docx
        doc = _docx.Document(str(path))

        paragraphs = []
        for para in doc.paragraphs:
            if para.text.strip():
                paragraphs.append(para.text)

        # Extraer tablas
        table_count = len(doc.tables)

        full_text = "\n\n".join(paragraphs)
        markdown_lines = []
        for para in doc.paragraphs:
            if not para.text.strip():
                continue
            style = para.style.name if para.style else ""
            if "Heading 1" in style:
                markdown_lines.append(f"# {para.text}")
            elif "Heading 2" in style:
                markdown_lines.append(f"## {para.text}")
            elif "Heading 3" in style:
                markdown_lines.append(f"### {para.text}")
            else:
                markdown_lines.append(para.text)
        markdown = "\n\n".join(markdown_lines)

        title = doc.core_properties.title if doc.core_properties else None
        word_count = len(full_text.split())

        return ParsedDocument(
            document_id=document_id,
            title=title or path.stem,
            plain_text=full_text,
            markdown=markdown,
            word_count=word_count,
            table_count=table_count,
            quality_score=0.8 if full_text else 0.1,
            parser_used="python-docx",
            parsed_at=now,
        )
    except Exception as exc:
        logger.warning("parse_docx %s: %s", path, exc)
        return ParsedDocument(
            document_id=document_id,
            parser_used="docx_fallback",
            warnings=[f"python-docx error: {exc}"],
            parsed_at=now,
        )


def parse_xlsx(path: str | Path) -> tuple["ParsedDocument", list["ExtractedTable"]]:
    """
    Parsea un archivo XLSX.

    Returns:
        (ParsedDocument, list[ExtractedTable])
    """
    from etl.sources.documents.schemas import ParsedDocument, ExtractedTable
    path = Path(path)
    document_id = path.stem
    now = datetime.now(timezone.utc)

    if not _PANDAS_OK and not _OPENPYXL_OK:
        return ParsedDocument(
            document_id=document_id,
            parser_used="xlsx_fallback",
            warnings=["pandas/openpyxl no instalados. Instala: pip install pandas openpyxl"],
            parsed_at=now,
        ), []

    tables = []
    sheets_text = []

    try:
        import pandas as pd
        xl = pd.ExcelFile(str(path))

        for sheet_idx, sheet_name in enumerate(xl.sheet_names):
            try:
                df = xl.parse(sheet_name)
                df_json = df.head(500).to_dict(orient="records")
                markdown = _df_to_markdown(df)
                sheets_text.append(f"## {sheet_name}\n\n{df.to_string()[:2000]}")

                tables.append(ExtractedTable(
                    document_id=document_id,
                    table_index=sheet_idx,
                    title=sheet_name,
                    dataframe_json={"records": df_json, "columns": list(df.columns)},
                    markdown=markdown,
                    extraction_method="pandas_xlsx",
                    confidence=0.9,
                ))
            except Exception as exc:
                logger.debug("parse_xlsx sheet %s: %s", sheet_name, exc)

        full_text = "\n\n".join(sheets_text)
        return ParsedDocument(
            document_id=document_id,
            title=path.stem,
            plain_text=full_text,
            markdown=full_text,
            table_count=len(tables),
            quality_score=0.85 if tables else 0.3,
            parser_used="pandas_xlsx",
            parsed_at=now,
        ), tables

    except Exception as exc:
        logger.warning("parse_xlsx %s: %s", path, exc)
        return ParsedDocument(
            document_id=document_id,
            parser_used="xlsx_fallback",
            warnings=[f"xlsx error: {exc}"],
            parsed_at=now,
        ), []


def parse_pptx(path: str | Path) -> "ParsedDocument":
    """
    Parsea un archivo PPTX. Stub seguro si no hay librería.
    Extrae texto de diapositivas si python-pptx está disponible.
    """
    from etl.sources.documents.schemas import ParsedDocument
    path = Path(path)
    document_id = path.stem
    now = datetime.now(timezone.utc)

    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides_text = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slides_text.append(f"## Diapositiva {i}\n\n" + "\n".join(texts))

        full_text = "\n\n".join(slides_text)
        return ParsedDocument(
            document_id=document_id,
            title=path.stem,
            plain_text=full_text,
            markdown=full_text,
            page_count=len(prs.slides),
            word_count=len(full_text.split()),
            quality_score=0.7 if full_text else 0.2,
            parser_used="python-pptx",
            parsed_at=now,
        )
    except ImportError:
        return ParsedDocument(
            document_id=document_id,
            parser_used="pptx_stub",
            warnings=["python-pptx no instalado. Instala: pip install python-pptx"],
            parsed_at=now,
        )
    except Exception as exc:
        logger.warning("parse_pptx %s: %s", path, exc)
        return ParsedDocument(
            document_id=document_id,
            parser_used="pptx_stub",
            warnings=[f"pptx error: {exc}"],
            parsed_at=now,
        )


def _df_to_markdown(df) -> str:
    try:
        lines = ["| " + " | ".join(str(c) for c in df.columns) + " |"]
        lines.append("|" + " --- |" * len(df.columns))
        for _, row in df.head(30).iterrows():
            lines.append("| " + " | ".join(str(v) for v in row) + " |")
        return "\n".join(lines)
    except Exception:
        return ""
